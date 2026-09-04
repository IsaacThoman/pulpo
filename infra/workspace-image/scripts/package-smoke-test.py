"""Exercise bundled packages without downloads, using real artifact round trips."""
import tempfile
from pathlib import Path
import sys
import subprocess

import numpy as np
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from PIL import Image
from docx import Document
from pptx import Presentation
from pypdf import PdfReader
import pdfplumber
from reportlab.pdfgen import canvas

assert sys.prefix == "/opt/pulpo/python", sys.prefix
with tempfile.TemporaryDirectory() as directory:
    root = Path(directory)
    frame = pd.DataFrame({"value": np.array([1, 2, 3])})
    frame.to_excel(root / "data.xlsx", index=False)
    assert pd.read_excel(root / "data.xlsx")["value"].sum() == 6
    plt.plot(frame["value"])
    plt.savefig(root / "plot.png")
    plt.close()
    with Image.open(root / "plot.png") as image:
        assert image.width > 0
    doc = Document()
    doc.add_paragraph("Pulpo document")
    doc.save(root / "document.docx")
    assert Document(root / "document.docx").paragraphs[0].text == "Pulpo document"
    deck = Presentation()
    deck.slides.add_slide(deck.slide_layouts[0])
    deck.save(root / "slides.pptx")
    assert len(Presentation(root / "slides.pptx").slides) == 1
    pdf = canvas.Canvas(str(root / "document.pdf"))
    pdf.drawString(72, 720, "Pulpo PDF")
    pdf.save()
    assert "Pulpo PDF" in PdfReader(root / "document.pdf").pages[0].extract_text()
    with pdfplumber.open(root / "document.pdf") as document:
        assert "Pulpo PDF" in document.pages[0].extract_text()
    subprocess.run(["pdftoppm", "-singlefile", "-scale-to", "100", "-png", str(root / "document.pdf"), str(root / "page")], check=True)
    with Image.open(root / "page.png") as image:
        assert max(image.size) == 100
print("Bundled package artifact tests passed.")
